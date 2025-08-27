import os
import zipfile
import json
import base64
import requests
from flask import Flask, request, jsonify, render_template, send_file
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import hashlib

# Configuration and setup
app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
GENERATED_FOLDER = 'generated'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(GENERATED_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['GENERATED_FOLDER'] = GENERATED_FOLDER

def get_slide_layout_by_name(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[0] # Fallback to first layout

def extract_images_from_template(pptx_path, temp_dir):
    image_paths = []
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
            for file_name in zip_ref.namelist():
                if file_name.startswith('ppt/media/') and (file_name.endswith('.jpeg') or file_name.endswith('.png')):
                    zip_ref.extract(file_name, temp_dir)
                    image_paths.append(os.path.join(temp_dir, file_name))
    except zipfile.BadZipFile:
        print("Not a valid PowerPoint file.")
    return image_paths

def generate_presentation_content(input_text, api_key, guidance, provider):
    """
    Calls the LLM API to generate presentation content.
    Returns a dictionary of slide data.
    """
    provider_endpoints = {
        'openai': 'https://api.openai.com/v1/chat/completions',
        'anthropic': 'https://api.anthropic.com/v1/messages',
        'gemini': 'https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent',
    }
    
    # Common instruction for all LLMs
    instruction = (
        f"You are a presentation content generator. Your task is to take a body of text and structure it "
        f"into a logical, multi-slide presentation. Each slide should have a title and a bulleted list of "
        f"key points. The content should be concise and easy to understand. "
        f"The user has provided the following guidance: '{guidance}'.\n\n"
        f"The output must be a JSON array of objects, where each object represents a slide. "
        f"Example format:\n"
        f"[{{"
        f"  \"title\": \"Slide Title\", "
        f"  \"points\": [\"Point 1\", \"Point 2\", \"Point 3\"]"
        f"}}, "
        f"{{"
        f"  \"title\": \"Another Slide\", "
        f"  \"points\": [\"Another point\"]"
        f"}}]\n\n"
        f"Now, please process the following text:\n\n---\n{input_text}\n---"
    )

    headers = {}
    payload = {}
    
    if provider == 'openai':
        headers['Authorization'] = f'Bearer {api_key}'
        payload = {
            "model": "gpt-4o-mini",
            "messages": [
                {"role": "system", "content": "You are a helpful assistant that generates presentation content in JSON format."},
                {"role": "user", "content": instruction}
            ]
        }
    elif provider == 'anthropic':
        headers = {
            'x-api-key': api_key,
            'anthropic-version': '2023-06-01',
            'content-type': 'application/json'
        }
        payload = {
            "model": "claude-3-haiku-20240307",
            "system": "You are a helpful assistant that generates presentation content in JSON format.",
            "messages": [
                {"role": "user", "content": instruction}
            ]
        }
    elif provider == 'gemini':
        headers = {
            'Content-Type': 'application/json',
        }
        url = f"{provider_endpoints[provider]}?key={api_key}"
        payload = {
            "contents": [
                {
                    "role": "user",
                    "parts": [
                        {"text": instruction}
                    ]
                }
            ]
        }
    else:
        return {'error': 'Unsupported LLM provider.'}, 400

    try:
        if provider == 'gemini':
            response = requests.post(url, headers=headers, json=payload)
        else:
            response = requests.post(provider_endpoints[provider], headers=headers, json=payload)
        response.raise_for_status()
        
        content = response.json()
        
        if provider == 'openai':
            json_str = content['choices'][0]['message']['content']
        elif provider == 'anthropic':
            json_str = content['content'][0]['text']
        elif provider == 'gemini':
            json_str = content['candidates'][0]['content']['parts'][0]['text']

        # Clean up markdown code blocks if present
        if json_str.startswith("```json"):
            json_str = json_str.strip("```json\n").strip()
            
        print("LLM JSON response:", json_str) # Log for debugging
        return json.loads(json_str)
        
    except requests.exceptions.RequestException as e:
        print(f"API call failed: {e}")
        return {'error': f"LLM API call failed: {e}"}, 500
    except json.JSONDecodeError:
        return {'error': 'LLM returned malformed JSON.'}, 500
    except Exception as e:
        print(f"An error occurred: {e}")
        return {'error': f"An error occurred: {e}"}, 500
    
def normalize_slide_data(raw):
    """Ensure slide data is always a list of dicts with 'title' and 'points'."""
    slides = []

    # If the root is a dict, wrap it in a list
    if isinstance(raw, dict):
        raw = [raw]

    # If the root is not a list, force a single slide
    if not isinstance(raw, list):
        raw = [{"title": str(raw), "points": []}]

    # Normalize each entry
    for i, item in enumerate(raw):
        if isinstance(item, dict):
            title = item.get("title", f"Slide {i+1}")
            points = item.get("points", [])
            if not isinstance(points, list):
                points = [str(points)]
            slides.append({"title": title, "points": points})
        else:
            slides.append({"title": str(item), "points": []})

    return slides


@app.route('/generate', methods=['POST'])
def generate_pptx():
    try:
        # Get data and file
        api_key = request.form.get('api_key')
        text_content = request.form.get('text_content')
        guidance = request.form.get('guidance', 'general purpose presentation')
        llm_provider = request.form.get('llm_provider', 'openai')
        
        if 'template' not in request.files:
            return jsonify({'error': 'No template file provided.'}), 400
        template_file = request.files['template']
        
        if not all([api_key, text_content, template_file]):
            return jsonify({'error': 'Missing required fields.'}), 400

        # Save template file
        template_path = os.path.join(app.config['UPLOAD_FOLDER'], template_file.filename)
        template_file.save(template_path)
        
        # Call LLM to get content
        slide_data = generate_presentation_content(text_content, api_key, guidance, llm_provider)
        if 'error' in slide_data:
            return jsonify(slide_data), slide_data.get('code', 500)

        # 🔧 Normalize the data to avoid 'int' or malformed structures
        slide_data = normalize_slide_data(slide_data)


        # Load template presentation
        try:
            prs = Presentation(template_path)
        except Exception as e:
            return jsonify({'error': f'Failed to load PowerPoint template: {e}'}), 500
        
        # Layouts (fallback if missing)
        title_slide_layout = prs.slide_layouts[0] if prs.slide_layouts else None
        content_slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else None
        
        # Fallback font and color
        template_font = "Calibri"
        from pptx.dml.color import RGBColor
        template_color = RGBColor(0, 0, 0)  # black

        # Remove existing slides
        for i in range(len(prs.slides) - 1, -1, -1):
            s = prs.slides._sldIdLst[i]
            prs.slides._sldIdLst.remove(s)
            
        # --- Title slide ---
        if title_slide_layout:
            title_slide = prs.slides.add_slide(title_slide_layout)
        else:
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])

        slide_title = slide_data[0].get("title", "Untitled Slide")
        title_shape = title_slide.shapes.title
        if title_shape:
            title_shape.text = slide_title
        else:
            textbox = title_slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            textbox.text_frame.text = slide_title

        # Add subtitle if available
        if slide_data[0].get("points"):
            if len(title_slide.placeholders) > 1:
                try:
                    subtitle_shape = title_slide.placeholders[1]
                    subtitle_shape.text = slide_data[0]["points"][0]
                except Exception:
                    pass  # skip if subtitle placeholder invalid

        # --- Other slides ---
        for i, slide in enumerate(slide_data[1:], start=1):
            layout = content_slide_layout or prs.slide_layouts[0]
            new_slide = prs.slides.add_slide(layout)
            
            # Title
            slide_title = slide.get("title", f"Slide {i+1}")
            title_shape = new_slide.shapes.title
            if title_shape:
                title_shape.text = slide_title
            else:
                textbox = new_slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                textbox.text_frame.text = slide_title
            
            # Bullet points
            if slide.get("points"):
                try:
                    body_shape = new_slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear()
                    for point in slide["points"]:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                except Exception:
                    # fallback: add textbox manually
                    left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4)
                    textbox = new_slide.shapes.add_textbox(left, top, width, height)
                    tf = textbox.text_frame
                    for point in slide["points"]:
                        p = tf.add_paragraph()
                        p.text = point

            # Apply font/color styling (fallbacks)
            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if template_font:
                                run.font.name = template_font
                            if template_color:
                                run.font.color.rgb = template_color

        # Save final PPTX
        output_filename = f"generated_{hashlib.md5(text_content.encode()).hexdigest()}.pptx"
        output_path = os.path.join(app.config['GENERATED_FOLDER'], output_filename)
        prs.save(output_path)
        
        return send_file(output_path, as_attachment=True, download_name='presentation.pptx')
        
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
        return jsonify({'error': f'An unexpected error occurred: {e}'}), 500

@app.route('/')
def index():
    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)